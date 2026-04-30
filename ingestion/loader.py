"""
Document loader:

Reads PDF, PPTX, TXT, and MD files and returns a list of Document objects,
each representing one logical page/slide with its source metadata attached.

Supported formats:
    .pdf  → one Document per page   (via PyMuPDF)
    .pptx → one Document per slide  (via python-pptx)
    .txt  → one Document for the whole file
    .md   → one Document for the whole file
"""

from __future__ import annotations

import logging
from dataclasses import dataclass, field
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt", ".md"}


@dataclass
class Document:
    """A single unit of raw text with its provenance metadata."""

    text: str
    metadata: dict = field(default_factory=dict)
    # metadata keys: source_file, page_number, total_pages, file_type


# ---------------------------------------------------------------------------
# Per-format readers
# ---------------------------------------------------------------------------

def _load_pdf(path: Path) -> list[Document]:
    """Extract text from each page of a PDF."""
    docs = []
    pdf = fitz.open(path)
    total = len(pdf)

    for i, page in enumerate(pdf):
        text = page.get_text("text").strip()
        if not text:
            logger.debug("Skipping empty page %d in %s", i + 1, path.name)
            continue
        docs.append(Document(
            text=text,
            metadata={
                "source_file": path.name,
                "page_number": i + 1,
                "total_pages": total,
                "file_type": "pdf",
            },
        ))

    pdf.close()
    return docs


def _load_pptx(path: Path) -> list[Document]:
    """Extract text from each slide of a PowerPoint file."""
    docs = []
    prs = Presentation(path)
    total = len(prs.slides)

    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        parts.append(line)

        text = "\n".join(parts)
        if not text:
            logger.debug("Skipping empty slide %d in %s", i + 1, path.name)
            continue

        docs.append(Document(
            text=text,
            metadata={
                "source_file": path.name,
                "page_number": i + 1,
                "total_pages": total,
                "file_type": "pptx",
            },
        ))

    return docs


def _load_text(path: Path) -> list[Document]:
    """Read a plain text or markdown file as a single document."""
    text = path.read_text(encoding="utf-8").strip()
    if not text:
        return []
    return [Document(
        text=text,
        metadata={
            "source_file": path.name,
            "page_number": 1,
            "total_pages": 1,
            "file_type": path.suffix.lstrip("."),
        },
    )]


# ---------------------------------------------------------------------------
# Public API -> rest of the app will call these
# ---------------------------------------------------------------------------

_LOADERS = {
    ".pdf": _load_pdf,
    ".pptx": _load_pptx,
    ".txt": _load_text,
    ".md": _load_text,
}


def load_file(path: str | Path) -> list[Document]:
    """
    Load a single file and return its Documents.

    Raises:
        ValueError: if the file type is not supported.
        FileNotFoundError: if the path doesn't exist.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    loader = _LOADERS.get(ext)
    if loader is None:
        raise ValueError(
            f"Unsupported file type '{ext}'. "
            f"Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )

    docs = loader(path)
    logger.info("Loaded %d document(s) from %s", len(docs), path.name)
    return docs


def load_directory(directory: str | Path) -> list[Document]:
    """
    Recursively load all supported files from a directory.

    Returns a flat list of Documents across all files.
    """
    directory = Path(directory)
    if not directory.is_dir():
        raise NotADirectoryError(f"Not a directory: {directory}")

    all_docs = []
    files = sorted(
        f for f in directory.rglob("*")
        if f.suffix.lower() in SUPPORTED_EXTENSIONS
    )

    if not files:
        logger.warning("No supported files found in %s", directory)
        return []

    for file_path in files:
        try:
            docs = load_file(file_path)
            all_docs.extend(docs)
        except Exception as e:
            logger.error("Failed to load %s: %s", file_path.name, e)

    logger.info(
        "Loaded %d document(s) from %d file(s) in %s",
        len(all_docs), len(files), directory,
    )
    return all_docs
